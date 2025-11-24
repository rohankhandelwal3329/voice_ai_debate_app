import io
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation


def extract_text_from_pdf(file_bytes: bytes) -> str:
    reader = PdfReader(io.BytesIO(file_bytes))
    pages = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            pages.append(text)
    return "\n".join(pages)


def extract_text_from_docx(file_bytes: bytes) -> str:
    doc = Document(io.BytesIO(file_bytes))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paragraphs)


def extract_text_from_pptx(file_bytes: bytes) -> str:
    prs = Presentation(io.BytesIO(file_bytes))
    slides_text = []
    for slide in prs.slides:
        slide_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_content.append(shape.text)
        if slide_content:
            slides_text.append(" ".join(slide_content))
    return "\n".join(slides_text)


def extract_text_from_txt(file_bytes: bytes) -> str:
    return file_bytes.decode("utf-8", errors="ignore")


def parse_file(filename: str, file_bytes: bytes) -> str:
    ext = filename.lower().split(".")[-1] if "." in filename else ""
    
    if ext == "pdf":
        return extract_text_from_pdf(file_bytes)
    elif ext == "docx":
        return extract_text_from_docx(file_bytes)
    elif ext == "pptx":
        return extract_text_from_pptx(file_bytes)
    elif ext == "txt":
        return extract_text_from_txt(file_bytes)
    else:
        raise ValueError(f"Unsupported file format: {ext}")

